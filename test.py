
import os

import hashlib
import asyncio
import concurrent.futures
from concurrent.futures import ThreadPoolExecutor
from docx2pdf import convert

from magic_pdf.data.data_reader_writer import FileBasedDataWriter, FileBasedDataReader
from magic_pdf.data.dataset import PymuDocDataset
from magic_pdf.model.doc_analyze_by_custom_model import doc_analyze
from magic_pdf.config.enums import SupportedPdfParseMethod

# 新增: 读取环境变量
from dotenv import load_dotenv

load_dotenv()


def calculate_file_hash(file_path, chunk_size=8192):
    """Calculate the hash value of a file."""
    hasher = hashlib.sha256()
    with open(file_path, "rb") as f:
        while chunk := f.read(chunk_size):
            hasher.update(chunk)
    return hasher.hexdigest()


def convert_ppt_to_pdf_alternative(file_path, temp_pdf):
    """
    Alternative method to convert PPT/PPTX to PDF using python-pptx and pypandoc.
    """
    try:
        import pypandoc
        # Convert PPT to intermediate format (e.g., HTML) using python-pptx
        intermediate_html = os.path.splitext(temp_pdf)[0] + ".html"
        # Export slides as images and combine into HTML (simplified example)
        from pptx import Presentation
        prs = Presentation(file_path)
        with open(intermediate_html, "w") as f:
            f.write("<html><body>")
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        f.write(f"<p>{shape.text}</p>")
            f.write("</body></html>")
        
        # Use pypandoc to convert HTML to PDF
        output = pypandoc.convert_file(intermediate_html, 'pdf', outputfile=temp_pdf)
        if output:
            raise RuntimeError(f"Error during conversion: {output}")
    except ImportError:
        raise RuntimeError("Alternative conversion requires pypandoc and python-pptx. Please install them.")


def convert_file_to_md(file_path, output_folder, image_writer, md_writer, image_dir):
    # 新增: 校验文件是否存在
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}. Skipping...")
        return

    # Get the file name and check if output exists
    name_without_suff = os.path.basename(file_path).split(".")[0]
    output_path = os.path.join(output_folder, f"{name_without_suff}.md")
    hash_path = os.path.join(output_folder, f"{name_without_suff}.hash")

    # Calculate current file hash
    current_hash = calculate_file_hash(file_path)

    # Check if hash file exists and compare hashes
    if os.path.exists(hash_path):
        with open(hash_path, "r") as hash_file:
            saved_hash = hash_file.read().strip()
        if saved_hash == current_hash:
            print(f"Skipping {file_path} - no changes detected")
            return

    try:  # 确保 try 块包含所有可能抛出异常的代码
        # Handle non-pdf files
        if file_path.lower().endswith(('.docx', 'doc', '.ppt', '.pptx')):
            # Convert to temporary pdf
            temp_pdf = os.path.join(os.path.dirname(file_path), "temp.pdf")
            
            if file_path.lower().endswith(('.docx', '.doc')):  # 修复: 支持多个后缀
                convert(file_path, temp_pdf)
            elif file_path.lower().endswith(('.ppt', '.pptx')):
                # 直接调用 alternative 方法
                convert_ppt_to_pdf_alternative(file_path, temp_pdf)

            # Process the temporary pdf
            file_path = temp_pdf

            # Read bytes from temporary pdf
            reader1 = FileBasedDataReader(os.path.dirname(file_path))
            pdf_bytes = reader1.read(os.path.basename(file_path))

            # Process the rest as before
            ds = PymuDocDataset(pdf_bytes)

            if ds.classify() == SupportedPdfParseMethod.OCR:
                infer_result = ds.apply(doc_analyze, ocr=True)
                pipe_result = infer_result.pipe_ocr_mode(image_writer)
            else:
                infer_result = ds.apply(doc_analyze, ocr=False)
                pipe_result = infer_result.pipe_txt_mode(image_writer)

            # Get markdown content
            md_content = pipe_result.get_markdown(image_dir)

            # Dump markdown
            pipe_result.dump_md(md_writer, f"{name_without_suff}.md", image_dir)
            
        else:
            # Original processing for pdf files
            reader1 = FileBasedDataReader(os.path.dirname(file_path))
            pdf_bytes = reader1.read(os.path.basename(file_path))  # read the pdf content

            # Create Dataset Instance
            ds = PymuDocDataset(pdf_bytes)

            # Inference
            if ds.classify() == SupportedPdfParseMethod.OCR:
                infer_result = ds.apply(doc_analyze, ocr=True)
                pipe_result = infer_result.pipe_ocr_mode(image_writer)
            else:
                infer_result = ds.apply(doc_analyze, ocr=False)
                pipe_result = infer_result.pipe_txt_mode(image_writer)

            # Get markdown content
            md_content = pipe_result.get_markdown(image_dir)

            # Dump markdown
            pipe_result.dump_md(md_writer, f"{name_without_suff}.md", image_dir)

    except Exception as e:  # 修复: 确保 except 匹配到正确的 try
        print(f"Error processing file {file_path}: {e}")

    finally:
        # Clean up temporary pdf
        temp_pdf = os.path.join(os.path.dirname(file_path), "temp.pdf")
        if os.path.exists(temp_pdf):
            os.remove(temp_pdf)

    # Save the current hash to the hash file
    with open(hash_path, "w") as hash_file:
        hash_file.write(current_hash)


# Ensure async_convert_pdfs_to_md is defined as a regular function
def async_convert_pdfs_to_md(input_folder, output_folder, num_threads):
    local_image_dir = os.path.join(output_folder, "images")
    local_md_dir = output_folder
    image_dir = str(os.path.basename(local_image_dir))

    os.makedirs(local_image_dir, exist_ok=True)

    image_writer = FileBasedDataWriter(local_image_dir)
    md_writer = FileBasedDataWriter(local_md_dir)

    # 修改: 支持更多文件类型
    file_paths = [
        os.path.join(input_folder, f)
        for f in os.listdir(input_folder)
        if f.lower().endswith(('.pdf', '.docx', '.doc', '.ppt', '.pptx'))
    ]

    # 新增: 打印文件路径，方便调试
    if not file_paths:
        print(f"No valid files found in the input folder: {input_folder}")
    else:
        print(f"Processing files: {file_paths}")

    with ThreadPoolExecutor(max_workers=num_threads) as executor:
        # 提交任务到线程池
        futures = [
            executor.submit(convert_file_to_md, file_path, output_folder, image_writer, md_writer, image_dir)
            for file_path in file_paths
        ]
        # 等待所有任务完成
        concurrent.futures.wait(futures)


if __name__ == "__main__":
    num_threads = 1  # 线程数量

    # 修改: 增加日志输出，确认环境变量是否正确加载
    input_folder = os.getenv('INPUT_FOLDER', './pdfs')
    output_folder = os.getenv('OUTPUT_FOLDER', './md')
    print(f"Input folder: {input_folder}")
    print(f"Output folder: {output_folder}")

    async_convert_pdfs_to_md(input_folder, output_folder, num_threads)
